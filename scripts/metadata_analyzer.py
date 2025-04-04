import os
import mimetypes
import exifread
from openpyxl.reader.excel import load_workbook
from pymediainfo import MediaInfo
from mutagen import File
from hachoir.parser import createParser
from hachoir.metadata import extractMetadata
from PyPDF2 import PdfReader
from docx import Document
import zipfile
from PIL import Image, ExifTags
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
from striprtf.striprtf import rtf_to_text
from io import BytesIO
from datetime import datetime
from pptx import Presentation
import json


def create_report_file(report_file_name):
    global folder_path
    folder_path = "C:/METKA_Reports"
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)
    global report_file_path
    report_file_path = os.path.join(folder_path, report_file_name)


def print_report_header(title):
    """Печатает заголовок отчета с рамкой '='"""
    print("=" * 50)
    print(f"{title.center(50)}")
    print("=" * 50)


def print_report_footer():
    """Печатает нижнюю строку отчета с рамкой '='"""
    print("=" * 50)
    print("Analysis completed | Анализ завершён.")
    print(f"The report is saved in | Отчёт сохранён в: {folder_path}")


def get_file_type(file_path):
    """Определяет MIME-тип файла, основываясь на его расширении."""
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type


def analyze_image_metadata(file_path):
    """Анализ метаданных изображений (EXIF)."""
    print_report_header("Image Analysis (EXIF) | Анализ изображения (EXIF)")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with open(file_path, 'rb') as f:
                tags = exifread.process_file(f)
                exif_tags = {tag: value for tag, value in tags.items() if not tag.startswith("JPEGThumbnail")}

                latitude = tags.get('GPSLatitude', None)
                longitude = tags.get('GPSLongitude', None)
                author = tags.get('Image Artist', None)

                if exif_tags:
                    report_file.write("EXIF metadata | EXIF-метаданные:\n")
                    print("EXIF metadata | EXIF-метаданные:")
                    for tag, value in exif_tags.items():
                        print(f"{tag}: {value}")
                        report_file.write(f"{tag}: {value}\n")

                if latitude and longitude:
                    report_file.write(f"Геолокация | Geolocation: {latitude}, {longitude}\n")
                    print(f"Геолокация | Geolocation: {latitude}, {longitude}")
                else:
                    report_file.write("Geolocation was not found | Геолокация не найдена.\n")
                    print("Geolocation was not found | Геолокация не найдена.")

                if author:
                    report_file.write(f"Author | Автор: {author}\n")
                    print(f"Author | Автор: {author}")
                else:
                    report_file.write("The author is not specified | Автор не указан.")
                    print("The author is not specified | Автор не указан.")

        except Exception as e:
            report_file.write(f"Error during EXIF analysis | Ошибка при анализе EXIF: {e}")
            print(f"Error during EXIF analysis | Ошибка при анализе EXIF: {e}")


def analyze_media_metadata(file_path):
    """Анализ метаданных видео-, аудио- и изображений (MediaInfo + Mutagen + EXIF)."""
    print_report_header("Media file analysis | Анализ медиафайла")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        metadata = {"Author": None, "GPS": None, "Company": None}

        try:
            media_info = MediaInfo.parse(file_path)
            if media_info.tracks:
                report_file.write("Media Data | Медиаданные:\n")
                print("Media Data | Медиаданные:")
                for track in media_info.tracks:
                    for key, value in track.to_data().items():
                        if key in ["Performer", "Encoded_By", "Publisher"]:
                            metadata["Author"] = value
                        if key in ["Copyright", "Publisher", "Organization", "Software"]:
                            metadata["Company"] = value
                        report_file.write(f"{key}: {value}\n")
                        print(f"{key}: {value}")
            else:
                report_file.write("Media data was not found | Медиаданные не найдены.\n")
                print("Media data was not found | Медиаданные не найдены.")
        except Exception as e:
            report_file.write(f"MediaInfo error | Ошибка MediaInfo: {e}\n")
            print(f"MediaInfo error | Ошибка MediaInfo: {e}")

        try:
            audio = File(file_path)
            if audio and audio.tags:
                report_file.write("Audio metadata | Аудиометаданные:\n")
                print("Audio metadata | Аудиометаданные:")
                for key, value in audio.tags.items():
                    if key in ["TPE1", "TPE2"]:
                        metadata["Author"] = value.text[0] if isinstance(value, list) else value
                    if key in ["TPUB"]:
                        metadata["Company"] = value.text[0] if isinstance(value, list) else value
                    report_file.write(f"{key}: {value}\n")
                    print(f"{key}: {value}")
            else:
                report_file.write("Audio metadata was not found | Аудиометаданные не найдены.\n")
                print("Audio metadata was not found | Аудиометаданные не найдены.")
        except Exception as e:
            report_file.write(f"Mutagen error | Ошибка Mutagen: {e}\n")
            print(f"Mutagen error | Ошибка Mutagen: {e}")

        try:
            if file_path.lower().endswith(('.jpg', '.jpeg', '.tiff', '.png')):
                with Image.open(file_path) as img:
                    exif_data = img._getexif()
                    if exif_data:
                        for tag, value in exif_data.items():
                            tag_name = ExifTags.TAGS.get(tag, tag)
                            if tag_name == "Artist":
                                metadata["Author"] = value
                            elif tag_name in ["GPSLatitude", "GPSLongitude"]:
                                if "GPS" not in metadata or metadata["GPS"] is None:
                                    metadata["GPS"] = {}
                                metadata["GPS"][tag_name] = value
                            elif tag_name == "Copyright":
                                metadata["Company"] = value
                            report_file.write(f"{tag_name}: {value}\n")
                            print(f"{tag_name}: {value}")
        except Exception as e:
            report_file.write(f"EXIF error | Ошибка EXIF: {e}\n")
            print(f"EXIF error | Ошибка EXIF: {e}")

        report_file.write(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}\n")
        report_file.write(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}\n")
        report_file.write(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}")

        print(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}")
        print(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}")
        print(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}")


def analyze_generic_metadata(file_path):
    """Анализ общих метаданных (Hachoir)."""
    print_report_header("General file analysis | Общий анализ файла")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        metadata_info = {"Author": None, "Company": None, "GPS": None}

        try:
            parser = createParser(file_path)
            if parser:
                metadata = extractMetadata(parser)
                if metadata:
                    report_file.write("Metadata | Метаданные:\n")
                    print("Metadata | Метаданные:")

                    for item in metadata.exportPlaintext():
                        report_file.write(item + "\n")
                        print(item)

                        lower_item = item.lower()
                        if "author" in lower_item or "created by" in lower_item:
                            metadata_info["Author"] = item.split(":", 1)[-1].strip()
                        if "company" in lower_item or "organization" in lower_item:
                            metadata_info["Company"] = item.split(":", 1)[-1].strip()
                        if "gps" in lower_item or "location" in lower_item:
                            metadata_info["GPS"] = item.split(":", 1)[-1].strip()

                else:
                    report_file.write("No metadata was found | Метаданные не найдены.\n")
                    print("No metadata was found | Метаданные не найдены.")
            else:
                report_file.write("Couldn't create a parser for Hachoir | Не удалось создать парсер для Hachoir.\n")
                print("Couldn't create a parser for Hachoir | Не удалось создать парсер для Hachoir.")
        except Exception as e:
            report_file.write(f"Error in Hachoir analysis | Ошибка при анализе Hachoir: {e}\n")
            print(f"Error in Hachoir analysis | Ошибка при анализе Hachoir: {e}")

        report_file.write(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}\n")
        report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
        report_file.write(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}\n")

        print(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}")
        print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
        print(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}")


def analyze_text_file(file_path):
    """Анализ текстового файла (кодировка, размер и метаданные)."""
    print_report_header("Text file analysis | Анализ текстового файла")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        metadata_info = {"Author": None, "Company": None, "GPS": None}

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                lines = content.splitlines()

                report_file.write(f"File Size | Размер файла: {os.path.getsize(file_path)} byte\n")
                report_file.write(
                    f"The first line of the file | Первая строка файла: {lines[0] if lines else 'Empty file | Пустой файл'}\n")
                report_file.write(f"Number of rows | Количество строк: {len(lines)}\n")

                print(f"File Size | Размер файла: {os.path.getsize(file_path)} byte")
                print(
                    f"The first line of the file | Первая строка файла: {lines[0] if lines else 'Empty file | Пустой файл'}")
                print(f"Number of rows | Количество строк: {len(lines)}")

                # Поиск автора, компании и геолокации в тексте
                for line in lines:
                    lower_line = line.lower()
                    if "author" in lower_line or "created by" in lower_line:
                        metadata_info["Author"] = line.split(":", 1)[-1].strip()
                    if "company" in lower_line or "organization" in lower_line:
                        metadata_info["Company"] = line.split(":", 1)[-1].strip()
                    if "gps" in lower_line or "location" in lower_line:
                        metadata_info["GPS"] = line.split(":", 1)[-1].strip()

            # Дата последней модификации файла
            last_modified_time = os.path.getmtime(file_path)
            report_file.write(
                f"Date of last modification | Дата последней модификации: {datetime.fromtimestamp(last_modified_time)}\n")
            print(
                f"Date of last modification | Дата последней модификации: {datetime.fromtimestamp(last_modified_time)}")

        except Exception as e:
            report_file.write(f"Error when analyzing a text file | Ошибка при анализе текстового файла: {e}\n")
            print(f"Error when analyzing a text file | Ошибка при анализе текстового файла: {e}")

        report_file.write(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}\n")
        report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
        report_file.write(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}\n")

        print(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}")
        print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
        print(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}")


def analyze_pdf_metadata(file_path):
    """Анализ метаданных PDF файла и встроенных изображений (EXIF)."""
    print_report_header("PDF file Analysis | Анализ PDF файла")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                metadata = reader.metadata
                report_file.write("Metadata of the PDF file | Метаданные PDF файла:\n")
                print("Metadata of the PDF file | Метаданные PDF файла:")
                for key, value in metadata.items():
                    report_file.write(f"{key}: {value}\n")
                    print(f"{key}: {value}")

        except Exception as e:
            report_file.write(f"Error in PDF analysis | Ошибка при анализе PDF: {e}")
            print(f"Error in PDF analysis | Ошибка при анализе PDF: {e}")


def analyze_pptx_metadata(file_path):
    """Анализ метаданных презентации PowerPoint и встроенных изображений (EXIF)."""
    print("PowerPoint Document Analysis | Анализ документа PowerPoint")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        metadata_info = {"Author": None, "Company": None, "GPS": None}

        try:
            prs = Presentation(file_path)
            props = prs.core_properties

            metadata_info["Author"] = props.author or "Not specified | Не указан"
            metadata_info["Company"] = getattr(props, 'company', 'Not specified | Не указана')

            report_file.write(f"Author | Автор: {metadata_info['Author']}\n")
            report_file.write(f"Company | Компания: {metadata_info['Company']}\n")
            report_file.write(f"Date of creation | Дата создания: {props.created or 'Unknown | Неизвестна'}\n")
            report_file.write(f"Date of last change | Дата последнего изменения: {props.modified or 'Unknown | Неизвестна'}\n")

            print(f"Author | Автор: {metadata_info['Author']}")
            print(f"Company | Компания: {metadata_info['Company']}")
            print(f"Date of creation | Дата создания: {props.created or 'Unknown | Неизвестна'}")
            print(f"Date of last change | Дата последнего изменения: {props.modified or 'Unknown | Неизвестна'}")

            slide_count = len(prs.slides)
            report_file.write(f"Number of slides | Количество слайдов: {slide_count}\n")
            print(f"Number of slides | Количество слайдов: {slide_count}")

            image_count = 0

            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if hasattr(shape, "image"):
                        image_count += 1
                        report_file.write(
                            f"\n\nImage {image_count} on the slide {slide_num} | Изображение {image_count} на слайде {slide_num}:\n")
                        print(
                            f"\nImage {image_count} on the slide {slide_num} | Изображение {image_count} на слайде {slide_num}:")

                        img_data = shape.image.blob  # Извлекаем бинарные данные изображения

                        try:
                            img = Image.open(BytesIO(img_data))
                            exif_data = img.getexif()

                            if exif_data:
                                report_file.write("EXIF image data | EXIF-данные изображения:\n")
                                print("EXIF image data | EXIF-данные изображения:")

                                for tag, value in exif_data.items():
                                    tag_name = ExifTags.TAGS.get(tag, tag)
                                    if tag_name in ["GPSLatitude", "GPSLongitude"]:
                                        if metadata_info["GPS"] is None:
                                            metadata_info["GPS"] = {}
                                        metadata_info["GPS"][tag_name] = value
                                    report_file.write(f"{tag_name}: {value}\n")
                                    print(f"{tag_name}: {value}")
                                    report_file.write(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}\n")
                                    print(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}")
                            else:
                                report_file.write(
                                    "\nNo EXIF data was found for the image | EXIF-данные для изображения не найдены.\n")
                                print("No EXIF data was found for the image | EXIF-данные для изображения не найдены.")
                        except Exception as e:
                            report_file.write(f"Error in image processing | Ошибка при обработке изображения: {e}\n")
                            print(f"Error in image processing | Ошибка при обработке изображения: {e}")

        except Exception as e:
            report_file.write(f"Error analyzing PowerPoint document | Ошибка при анализе PowerPoint документа: {e}\n")
            print(f"Error analyzing PowerPoint document | Ошибка при анализе PowerPoint документа: {e}")


def analyze_word_metadata(file_path):
    """Анализ метаданных документа Word и встроенных изображений (EXIF)."""
    print_report_header("Word Document Analysis | Анализ документа Word")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        metadata_info = {"Author": None, "Company": None, "GPS": None}

        try:
            doc = Document(file_path)
            core_props = doc.core_properties

            metadata_info["Author"] = core_props.author or "Not specified | Не указан"
            metadata_info["Company"] = getattr(core_props, 'company', 'Not specified | Не указана')

            report_file.write(f"Author | Автор: {metadata_info['Author']}\n")
            report_file.write(f"Company | Компания: {metadata_info['Company']}\n")
            report_file.write(f"Date of creation | Дата создания: {core_props.created or 'Unknown | Неизвестна'}\n")
            report_file.write(f"Date of last change | Дата последнего изменения: {core_props.modified or 'Unknown | Неизвестна'}\n")

            print(f"Author | Автор: {metadata_info['Author']}")
            print(f"Company | Компания: {metadata_info['Company']}")
            print(f"Date of creation | Дата создания: {core_props.created or 'Unknown | Неизвестна'}")
            print(f"Date of last change | Дата последнего изменения: {core_props.modified or 'Unknown | Неизвестна'}")

            report_file.write(f"Number of paragraphs | Количество параграфов: {len(doc.paragraphs)}\n")
            report_file.write(f"Number of tables | Количество таблиц: {len(doc.tables)}\n")
            print(f"Number of paragraphs | Количество параграфов: {len(doc.paragraphs)}")
            print(f"Number of tables | Количество таблиц: {len(doc.tables)}")

            image_count = 0

            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    image_count += 1
                    report_file.write(f"\n\nImage {image_count} in document | Изображение {image_count} в документе:\n")
                    print(f"\nImage {image_count} in document | Изображение {image_count} в документе:")

                    img_data = rel.target_part.blob

                    try:
                        img = Image.open(BytesIO(img_data))
                        exif_data = img.getexif()

                        if exif_data:
                            report_file.write("EXIF image data | EXIF-данные изображения:\n")
                            print("EXIF image data | EXIF-данные изображения:")

                            for tag, value in exif_data.items():
                                tag_name = ExifTags.TAGS.get(tag, tag)
                                if tag_name in ["GPSLatitude", "GPSLongitude"]:
                                    if metadata_info["GPS"] is None:
                                        metadata_info["GPS"] = {}
                                    metadata_info["GPS"][tag_name] = value
                                report_file.write(f"{tag_name}: {value}\n")
                                print(f"{tag_name}: {value}")
                                report_file.write(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}\n")
                                print(f"Geolocation | Геолокация: {metadata_info['GPS'] or 'Not specified | Не указана'}")
                        else:
                            report_file.write(
                                "No EXIF data was found for the image | EXIF-данные для изображения не найдены.\n")
                            print("No EXIF data was found for the image | EXIF-данные для изображения не найдены.")
                    except Exception as e:
                        report_file.write(f"Error in image processing | Ошибка при обработке изображения: {e}\n")
                        print(f"Error in image processing | Ошибка при обработке изображения: {e}")

        except Exception as e:
            report_file.write(f"Error when analyzing a Word document | Ошибка при анализе Word документа: {e}\n")
            print(f"Error when analyzing a Word document | Ошибка при анализе Word документа: {e}")


def analyze_excel_metadata(file_path):
    """Анализ метаданных документа Excel и встроенных изображений (EXIF)."""
    print_report_header("Excel Document Analysis | Анализ документа Excel")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            wb = load_workbook(file_path)
            props = wb.properties
            report_file.write(f"Author | Автор: {props.creator or 'Not specified | Не указан'}\n")
            print(f"Author | Автор: {props.creator or 'Not specified | Не указан'}")
            company = getattr(props, 'company', 'Не указана')
            report_file.write(f"Company | Компания: {company}\n")
            print(f"Company | Компания: {company}")
            report_file.write(f"Date of creation | Дата создания: {props.created or 'Unknown | Неизвестна'}\n")
            report_file.write(
                f"Date of last change | Дата последнего изменения: {props.modified or 'Unknown | Неизвестна'}\n")
            print(f"Date of creation | Дата создания: {props.created or 'Unknown | Неизвестна'}")
            print(f"Date of last change | Дата последнего изменения: {props.modified or 'Unknown | Неизвестна'}")

            sheet_names = wb.sheetnames
            report_file.write(f"Number of sheets | Количество листов: {len(sheet_names)}\n")
            report_file.write(f"Sheet Names | Имена листов: {', '.join(sheet_names)}\n")
            print(f"Number of sheets | Количество листов: {len(sheet_names)}")
            print(f"Sheet Names | Имена листов: {', '.join(sheet_names)}")

            image_count = 0
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws._images:
                    image_count += 1
                    img_data = row._data()
                    try:
                        img = Image.open(BytesIO(img_data))
                        exif_data = img.getexif()
                        if exif_data:
                            report_file.write(
                                f"\nEXIF image data from sheet {sheet} | EXIF-данные изображения с листа {sheet}:\n")
                            print(f"\nEXIF image data from sheet {sheet} | EXIF-данные изображения с листа {sheet}:")
                            for tag, value in exif_data.items():
                                tag_name = ExifTags.TAGS.get(tag, tag)
                                report_file.write(f"{tag_name}: {value}\n")
                                print(f"{tag_name}: {value}")
                                if tag_name in ["GPSLatitude", "GPSLongitude"]:
                                    report_file.write(f"Geolocation | Геолокация: {value}\n")
                                    print(f"Geolocation | Геолокация: {value}")
                        else:
                            report_file.write(
                                "No EXIF data was found for the image | EXIF-данные для изображения не найдены.\n")
                            print("No EXIF data was found for the image | EXIF-данные для изображения не найдены.")
                    except Exception as e:
                        report_file.write(f"Error in image processing | Ошибка при обработке изображения: {e}\n")
                        print(f"Error in image processing | Ошибка при обработке изображения: {e}")
        except Exception as e:
            report_file.write(f"Error when analyzing an Excel document | Ошибка при анализе Excel документа: {e}")
            print(f"Error when analyzing an Excel document | Ошибка при анализе Excel документа: {e}")


def analyze_zip_metadata(file_path):
    """Анализ метаданных ZIP архива и встроенных изображений (EXIF)."""
    print_report_header("ZIP Archive Analysis | Анализ ZIP архива")

    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                report_file.write(f"Number of files in the archive | Количество файлов в архиве: {len(zip_ref.namelist())}")
                print(f"Number of files in the archive | Количество файлов в архиве: {len(zip_ref.namelist())}")
                for file in zip_ref.namelist():
                    report_file.write(f"\nThe file is in the archive | Файл в архиве: {file}\n")
                    print(f"\nThe file is in the archive | Файл в архиве: {file}")

                    if file.lower().endswith(('.jpg', '.jpeg', '.png', '.tiff', '.bmp')):
                        with zip_ref.open(file) as img_file:
                            img_data = img_file.read()

                            try:
                                img = Image.open(BytesIO(img_data))
                                exif_data = img.getexif()

                                if exif_data:
                                    report_file.write(f"EXIF image data | EXIF-данные изображения {file}:\n")
                                    print(f"EXIF image data | EXIF-данные изображения {file}:")
                                    for tag, value in exif_data.items():
                                        report_file.write(f"{tag}: {value}\n")
                                        print(f"{tag}: {value}")
                                else:
                                    report_file.write(f"The EXIF data for the {file} image was not found | EXIF-данные для изображения {file} не найдены.\n")
                                    print(f"The EXIF data for the {file} image was not found | EXIF-данные для изображения {file} не найдены.")
                            except Exception as e:
                                report_file.write(f"Error in image processing {file} | Ошибка при обработке изображения {file}: {e}\n")
                                print(f"Error in image processing {file} | Ошибка при обработке изображения {file}: {e}")

        except Exception as e:
            report_file.write(f"Error when analyzing the ZIP archive | Ошибка при анализе ZIP архива: {e}")
            print(f"Error when analyzing the ZIP archive | Ошибка при анализе ZIP архива: {e}")


def analyze_tiff_image(file_path):
    """Анализ TIFF-изображения и EXIF-метаданных."""
    print_report_header("TIFF Image Analysis | Анализ TIFF изображения")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with Image.open(file_path) as img:
                report_file.write(f"Image Format | Формат изображения: {img.format}\n")
                report_file.write(f"Size | Размер: {img.size[0]}x{img.size[1]}\n")
                report_file.write(f"The color palette | Цветовая палитра: {img.mode}\n")
                print(f"Image Format | Формат изображения: {img.format}")
                print(f"Size | Размер: {img.size[0]}x{img.size[1]}")
                print(f"The color palette | Цветовая палитра: {img.mode}")

                exif_data = img.getexif()
                if exif_data:
                    report_file.write("EXIF image data | EXIF-данные изображения:\n")
                    print("EXIF image data | EXIF-данные изображения:")
                    metadata = {"Author": None, "Company": None, "GPS": None}
                    for tag, value in exif_data.items():
                        tag_name = ExifTags.TAGS.get(tag, tag)
                        if tag_name == "Artist":
                            metadata["Author"] = value
                        elif tag_name == "Make":
                            metadata["Company"] = value
                        elif tag_name == "GPSLatitude" or tag_name == "GPSLongitude":
                            if "GPS" not in metadata or metadata["GPS"] is None:
                                metadata["GPS"] = {}
                            metadata["GPS"][tag_name] = value
                        report_file.write(f"{tag_name}: {value}\n")
                        print(f"{tag_name}: {value}")

                    report_file.write(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}\n")
                    report_file.write(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}\n")
                    report_file.write(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}\n")
                    print(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}")
                    print(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}")
                    print(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}")
                else:
                    report_file.write("EXIF data was not found for the image | EXIF-данные не найдены для изображения.")
                    print("EXIF data was not found for the image | EXIF-данные не найдены для изображения.")
        except Exception as e:
            report_file.write(f"Error in TIFF image analysis | Ошибка при анализе TIFF изображения: {e}")
            print(f"Error in TIFF image analysis | Ошибка при анализе TIFF изображения: {e}")


def analyze_png_image(file_path):
    """Анализ PNG изображения."""
    print_report_header("PNG Image Analysis | Анализ PNG изображения")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with Image.open(file_path) as img:
                report_file.write(f"Image Format | Формат изображения: {img.format}\n")
                report_file.write(f"Size | Размер: {img.size[0]}x{img.size[1]}\n")
                report_file.write(f"The color palette | Цветовая палитра: {img.mode}\n")
                print(f"Image Format | Формат изображения: {img.format}")
                print(f"Size | Размер: {img.size[0]}x{img.size[1]}")
                print(f"The color palette | Цветовая палитра: {img.mode}")

                if img.info.get('text'):
                    report_file.write("\nText comments images | Текстовые комментарии изображения:\n")
                    print("\nText comments images | Текстовые комментарии изображения:")
                    for key, value in img.info['text'].items():
                        report_file.write(f"{key}: {value}\n")
                        print(f"{key}: {value}")
                else:
                    report_file.write("\nNo text comments found | Текстовые комментарии не найдены.\n")
                    print("\nNo text comments found | Текстовые комментарии не найдены.")

                if 'compression' in img.info:
                    report_file.write(f"Type of compression | Тип сжатия: {img.info['compression']}\n")
                    print(f"Type of compression | Тип сжатия: {img.info['compression']}")
                else:
                    report_file.write("Compression type not found | Тип сжатия не найден.\n")
                    print("Compression type not found | Тип сжатия не найден.")

                exif_data = img._getexif()
                if exif_data:
                    report_file.write("EXIF data | EXIF-данные:\n")
                    print("EXIF data | EXIF-данные:")
                    metadata = {"Author": None, "Company": None, "GPS": None}
                    for tag, value in exif_data.items():
                        tag_name = ExifTags.TAGS.get(tag, tag)
                        if tag_name == "Artist":
                            metadata["Author"] = value
                        elif tag_name == "Make":
                            metadata["Company"] = value
                        elif tag_name == "GPSLatitude" or tag_name == "GPSLongitude":
                            if "GPS" not in metadata or metadata["GPS"] is None:
                                metadata["GPS"] = {}
                            metadata["GPS"][tag_name] = value
                        report_file.write(f"{tag_name}: {value}\n")
                        print(f"{tag_name}: {value}")
                    report_file.write(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}\n")
                    report_file.write(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}\n")
                    report_file.write(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}\n")
                    print(f"Author | Автор: {metadata['Author'] or 'Not specified | Не указан'}")
                    print(f"Company | Компания: {metadata['Company'] or 'Not specified | Не указана'}")
                    print(f"Geolocation | Геолокация: {metadata['GPS'] or 'Not specified | Не указана'}")
                else:
                    report_file.write("EXIF data not found | EXIF-данные не найдены.")
                    print("EXIF data not found | EXIF-данные не найдены.")
        except Exception as e:
            report_file.write(f"Error analyzing PNG image | Ошибка при анализе PNG изображения: {e}")
            print(f"Error analyzing PNG image | Ошибка при анализе PNG изображения: {e}")


def analyze_xml(file_path):
    """Анализ XML документа."""
    print_report_header("XML document analysis | Анализ XML документа")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            report_file.write(f"The root element | Корневой элемент: {root.tag}\n")
            print(f"The root element | Корневой элемент: {root.tag}")
            metadata = root.attrib

            if metadata:
                report_file.write("Metadata of an XML document | Метаданные XML документа:\n")
                print("Metadata of an XML документа | Метаданные XML документа:")
                metadata_info = {"Author": None, "Company": None,
                                 "Geolocation": None}
                for key, value in metadata.items():
                    report_file.write(f"{key}: {value}\n")
                    print(f"{key}: {value}")

                    if key.lower() == "author":
                        metadata_info["Author"] = value
                    elif key.lower() == "company":
                        metadata_info["Company"] = value
                    elif key.lower() == "geolocation":
                        metadata_info["Geolocation"] = value

                report_file.write(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}\n")
                report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
                report_file.write(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}\n")
                print(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}")
                print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
                print(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}")
            else:
                report_file.write("Metadata not found | Метаданные не найдены.\n")
                print("Metadata not found | Метаданные не найдены.")

            report_file.write("XML Elements | Элементы XML:\n")
            print("XML Elements | Элементы XML:")
            for elem in root.iter():
                report_file.write(f"{elem.tag}: {elem.text}\n")
                print(f"{elem.tag}: {elem.text}")
        except Exception as e:
            report_file.write(f"Error when analyzing an XML document | Ошибка при анализе XML документа: {e}\n")
            print(f"Error when analyzing an XML document | Ошибка при анализе XML документа: {e}")


def analyze_html(file_path):
    """Анализ HTML документа."""
    print_report_header("HTML document analysis | Анализ HTML документа")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')

                report_file.write("Metadata of an HTML document | Метаданные HTML документа:\n")
                print("Metadata of an HTML document | Метаданные HTML документа:")
                meta_tags = soup.find_all('meta')
                metadata_info = {"Author": None, "Company": None, "Geolocation": None}
                for meta in meta_tags:
                    if 'name' in meta.attrs:
                        report_file.write(f"{meta.attrs.get('name')}: {meta.attrs.get('content')}\n")
                        print(f"{meta.attrs.get('name')}: {meta.attrs.get('content')}")

                        if meta.attrs.get('name').lower() == "author":
                            metadata_info["Author"] = meta.attrs.get('content')
                        elif meta.attrs.get('name').lower() == "company":
                            metadata_info["Company"] = meta.attrs.get('content')
                        elif meta.attrs.get('name').lower() == "geolocation":
                            metadata_info["Geolocation"] = meta.attrs.get('content')
                    elif 'property' in meta.attrs:
                        report_file.write(f"{meta.attrs.get('property')}: {meta.attrs.get('content')}\n")
                        print(f"{meta.attrs.get('property')}: {meta.attrs.get('content')}")

                report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
                report_file.write(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}\n")
                print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
                print(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}")

                report_file.write(f"The tag | Тег <title>: {soup.title.string if soup.title else 'Not found | Не найден'}\n")
                report_file.write(f"Number of tags | Количество тегов <p>: {len(soup.find_all('p'))}\n")
                print(f"The tag | Тег <title>: {soup.title.string if soup.title else 'Not found | Не найден'}")
                print(f"Number of tags | Количество тегов <p>: {len(soup.find_all('p'))}")
        except Exception as e:
            report_file.write(f"Error when analyzing an HTML document | Ошибка при анализе HTML документа: {e}\n")
            print(f"Error when analyzing an HTML document | Ошибка при анализе HTML документа: {e}")


def analyze_rtf(file_path):
    """Анализ RTF документа."""
    print_report_header("RTF document analysis | Анализ RTF документа")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                rtf_content = f.read()
                text = rtf_to_text(rtf_content)
                report_file.write(f"RTF file size | Размер RTF файла: {os.path.getsize(file_path)} bite\n")
                report_file.write(f"The total number of characters in the document | Общее количество символов в документе: {len(text)}\n")
                report_file.write(f"The number of lines in the document | Количество строк в документе: {text.count('\n')}\n")
                report_file.write(f"The number of words in the document | Количество слов в документе: {len(text.split())}\n")
                print(f"RTF file size | Размер RTF файла: {os.path.getsize(file_path)} bite")
                print(f"The total number of characters in the document | Общее количество символов в документе: {len(text)}")
                print(f"The number of lines in the document | Количество строк в документе: {text.count('\n')}")
                print(f"The number of words in the document | Количество слов в документе: {len(text.split())}")

                metadata_info = {"Author": None, "Company": None, "Geolocation": None}

                if '\\author' in rtf_content:
                    author_start = rtf_content.find('\\author') + len('\\author')
                    author_end = rtf_content.find('}', author_start)
                    metadata_info["Author"] = rtf_content[author_start:author_end].strip()
                if '\\company' in rtf_content:
                    company_start = rtf_content.find('\\company') + len('\\company')
                    company_end = rtf_content.find('}', company_start)
                    metadata_info["Company"] = rtf_content[company_start:company_end].strip()
                if '\\geolocation' in rtf_content:
                    geolocation_start = rtf_content.find('\\geolocation') + len('\\geolocation')
                    geolocation_end = rtf_content.find('}', geolocation_start)
                    metadata_info["Geolocation"] = rtf_content[geolocation_start:geolocation_end].strip()

                report_file.write(f"\nAuthor | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}\n")
                report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
                report_file.write(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}\n")
                print(f"\nAuthor | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}")
                print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
                print(f"Geolocation | Геолокация: {metadata_info['Geolocation'] or 'Not specified | Не указана'}")
        except Exception as e:
            report_file.write(f"Error when analyzing an RTF document | Ошибка при анализе RTF документа: {e}\n")
            print(f"Error when analyzing an RTF document | Ошибка при анализе RTF документа: {e}")


def analyze_gif_image(file_path):
    """Анализ GIF изображения."""
    print_report_header("GIF analysis | Анализ GIF")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with Image.open(file_path) as img:
                report_file.write(f"Format | Формат: {img.format}\n")
                report_file.write(f"Size | Размер: {img.size[0]}x{img.size[1]}\n")
                report_file.write(f"Number of frames | Количество кадров: {img.n_frames}\n")
                print(f"Format | Формат: {img.format}")
                print(f"Size | Размер: {img.size[0]}x{img.size[1]}")
                print(f"Number of frames | Количество кадров: {img.n_frames}")

        except Exception as e:
            report_file.write(f"Error when analyzing a GIF image | Ошибка при анализе GIF изображения: {e}\n")
            print(f"Error when analyzing a GIF image | Ошибка при анализе GIF изображения: {e}")


def analyze_json(file_path):
    """Анализ JSON-файла."""
    print_report_header("JSON Metadata Analysis | Анализ JSON метаданных")
    with open(report_file_path, "w", encoding="utf-8") as report_file:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                metadata = json.load(f)

            metadata_info = {
                "Author": None,
                "Company": None,
                "Geolocation": None
            }

            report_file.write("Extracted Metadata | Извлечённые метаданные:\n")
            print("Extracted Metadata | Извлечённые метаданные:")

            for key, value in metadata.items():
                report_file.write(f"{key}: {value}\n")
                print(f"{key}: {value}")
                lower_key = key.lower()
                if "author" in lower_key or "creator" in lower_key:
                    metadata_info["Author"] = value
                elif "company" in lower_key:
                    metadata_info["Company"] = value
                elif "gps" in lower_key:
                    if metadata_info["Geolocation"] is None:
                        metadata_info["Geolocation"] = []
                    metadata_info["Geolocation"].append((key, value))

            report_file.write(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}\n")
            report_file.write(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}\n")
            if metadata_info["Geolocation"]:
                geo_str = ", ".join(f"{k}: {v}" for k, v in metadata_info["Geolocation"])
            else:
                geo_str = "Not specified | Не указана"
            report_file.write(f"Geolocation | Геолокация: {geo_str}\n")
            print(f"Author | Автор: {metadata_info['Author'] or 'Not specified | Не указан'}")
            print(f"Company | Компания: {metadata_info['Company'] or 'Not specified | Не указана'}")
            print(f"Geolocation | Геолокация: {geo_str}")

        except Exception as e:
            error_msg = f"Error while analyzing JSON metadata | Ошибка при анализе JSON: {e}"
            report_file.write(error_msg + "\n")
            print(error_msg)


def analyze_metadata(file_path):
    """Определяет тип файла и запускает нужный метод анализа."""
    if not os.path.exists(file_path):
        print("The file was not found | Файл не найден.")
        return

    # Определяем тип файла
    mime_type = get_file_type(file_path)
    print(f"\nMIME-type: {mime_type if mime_type else "Couldn't determine the file type | Не удалось определить тип файла"}")

    # Проверяем расширение файла и корректируем MIME-тип
    if mime_type == "application/msword" and file_path.lower().endswith(".rtf"):
        mime_type = "application/rtf"

    # Анализ в зависимости от типа
    if mime_type:
        if mime_type.startswith("image"):
            if file_path.lower().endswith(".tiff"):
                analyze_tiff_image(file_path)
            elif file_path.lower().endswith(".png"):
                analyze_png_image(file_path)
            elif file_path.lower().endswith(".gif"):
                analyze_gif_image(file_path)
            else:
                analyze_image_metadata(file_path)
        elif mime_type.startswith(("audio", "video")):
            analyze_media_metadata(file_path)
        elif mime_type == "application/pdf":
            analyze_pdf_metadata(file_path)
        elif mime_type in [
            "application/msword",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ]:
            analyze_word_metadata(file_path)
        elif mime_type in [
            "application/vnd.ms-excel",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ]:
            analyze_excel_metadata(file_path)
        elif mime_type in [
            "application/vnd.ms-powerpoint",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ]:
            analyze_pptx_metadata(file_path)
        elif mime_type in [
            "application/x-zip-compressed",
            "application/zip"
        ]:
            analyze_zip_metadata(file_path)
        elif mime_type == "application/xml":
            analyze_xml(file_path)
        elif mime_type == "text/html":
            analyze_html(file_path)
        elif mime_type == "application/rtf":
            analyze_rtf(file_path)
        elif mime_type.startswith("text"):
            analyze_text_file(file_path)
        elif mime_type == "application/json":
            analyze_json(file_path)
        else:
            analyze_generic_metadata(file_path)
    else:
        analyze_generic_metadata(file_path)

    print_report_footer()